# financial_metrics_extractor.py
import os
import re
import json
import anthropic
from datetime import datetime, timedelta
from decimal import Decimal
import pandas as pd
import PyPDF2
from pptx import Presentation
import openpyxl
from io import BytesIO

from models import Company, EmailUpdate, Attachment
from financial_metrics_models import FinancialMetrics, MetricExtraction
from db import SessionLocal

class FinancialMetricsExtractor:
    def __init__(self):
        self.claude_api_key = os.getenv('ANTHROPIC_API_KEY')
        if not self.claude_api_key:
            raise ValueError("ANTHROPIC_API_KEY environment variable not set")
        
        self.claude_client = anthropic.Anthropic(api_key=self.claude_api_key)
        
        # Financial metrics we're looking for
        self.target_metrics = [
            'mrr', 'arr', 'qrr', 'total_revenue', 'gross_revenue', 'net_revenue',
            'mrr_growth', 'arr_growth', 'revenue_growth_yoy', 'revenue_growth_mom',
            'cash_balance', 'net_burn', 'gross_burn', 'runway_months',
            'gross_margin', 'ebitda', 'ebitda_margin', 'net_income',
            'customer_count', 'new_customers', 'churn_rate', 'ltv', 'cac',
            'team_size', 'bookings', 'pipeline'
        ]
        
        print(f"💰 Financial Metrics Extractor initialized")
        print(f"🤖 Claude API: {'✅ Connected' if self.claude_api_key else '❌ Not configured'}")

    def extract_text_from_pdf(self, file_path):
        """Extract text from PDF file"""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                return text
        except Exception as e:
            print(f"❌ Error extracting PDF text: {e}")
            return None

    def extract_text_from_excel(self, file_path):
        """Extract text and data from Excel file"""
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            extracted_data = {}
            
            for sheet_name in excel_file.sheet_names:
                try:
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    # Convert to string representation
                    sheet_text = f"Sheet: {sheet_name}\n"
                    sheet_text += df.to_string(max_cols=10, max_rows=50)
                    extracted_data[sheet_name] = sheet_text
                except Exception as e:
                    print(f"⚠️  Error reading sheet {sheet_name}: {e}")
                    continue
            
            # Combine all sheets
            full_text = "\n\n".join(extracted_data.values())
            return full_text
            
        except Exception as e:
            print(f"❌ Error extracting Excel data: {e}")
            return None

    def extract_text_from_pptx(self, file_path):
        """Extract text from PowerPoint file"""
        try:
            presentation = Presentation(file_path)
            text = ""
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                text += f"Slide {slide_num}:\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                
                text += "\n"
            
            return text
            
        except Exception as e:
            print(f"❌ Error extracting PowerPoint text: {e}")
            return None

    def extract_attachment_content(self, attachment):
        """Extract content from attachment based on file type"""
        file_path = attachment.path
        filename = attachment.filename.lower()
        
        print(f"📄 Extracting content from: {attachment.filename}")
        
        if filename.endswith('.pdf'):
            return self.extract_text_from_pdf(file_path)
        elif filename.endswith(('.xlsx', '.xls')):
            return self.extract_text_from_excel(file_path)
        elif filename.endswith(('.pptx', '.ppt')):
            return self.extract_text_from_pptx(file_path)
        else:
            print(f"⚠️  Unsupported file type: {filename}")
            return None

    def analyze_content_with_claude(self, content, company_name, source_info):
        """Use Claude to extract financial metrics from content"""
        
        prompt = f"""
You are a financial analyst extracting key metrics from portfolio company updates. 

Company: {company_name}
Source: {source_info}

Please analyze the following content and extract financial metrics. Return ONLY a JSON object with the following structure:

{{
    "reporting_period": "Q1 2025" or "May 2025" or "2024 Annual" etc,
    "reporting_date": "2025-05-01" (best estimate in YYYY-MM-DD format),
    "mrr": "Monthly Recurring Revenue (e.g., '$112K', '$1.2M') or 'N/A'",
    "arr": "Annual Recurring Revenue (e.g., '$8.022M', '~$8.000M') or 'N/A'", 
    "qrr": "Quarterly Recurring Revenue or 'N/A'",
    "total_revenue": "Total revenue for period or 'N/A'",
    "gross_revenue": "Gross revenue or 'N/A'",
    "net_revenue": "Net revenue or 'N/A'",
    "mrr_growth": "MRR growth rate (e.g., '+15%', '-5%') or 'N/A'",
    "arr_growth": "ARR growth rate (e.g., '+148.509%') or 'N/A'",
    "revenue_growth_yoy": "Year over year growth or 'N/A'",
    "revenue_growth_mom": "Month over month growth or 'N/A'",
    "cash_balance": "Current cash balance (e.g., '$2.8M', '$1.35M') or 'N/A'",
    "net_burn": "Monthly net burn rate or 'N/A'",
    "gross_burn": "Monthly gross burn rate or 'N/A'",
    "runway_months": "Cash runway in months (e.g., '24+ months') or 'N/A'",
    "gross_margin": "Gross margin percentage (e.g., '72%') or 'N/A'",
    "ebitda": "EBITDA or 'N/A'",
    "ebitda_margin": "EBITDA margin or 'N/A'",
    "net_income": "Net income/loss or 'N/A'",
    "customer_count": "Total customers (e.g., '50 clients') or 'N/A'",
    "new_customers": "New customers in period or 'N/A'",
    "churn_rate": "Customer churn rate or 'N/A'",
    "ltv": "Lifetime value or 'N/A'",
    "cac": "Customer acquisition cost or 'N/A'",
    "team_size": "Number of employees or 'N/A'",
    "bookings": "New bookings/contracts (e.g., '$42K in early revenue') or 'N/A'",
    "pipeline": "Sales pipeline value or 'N/A'",
    "key_highlights": "Key achievements and positive developments",
    "key_challenges": "Challenges and concerns mentioned",
    "funding_status": "Current funding status/notes or 'N/A'",
    "extraction_confidence": "high/medium/low based on clarity of data"
}}

Key guidelines:
- Preserve original formatting (e.g., "$1.2M", "~$8.000M", "24+ months")
- Use "N/A" for metrics not mentioned or unclear
- Be conservative with confidence - use "low" if uncertain
- Extract exact numbers and formatting as presented
- Look for terms like: MRR, ARR, revenue, cash, burn, runway, customers, etc.

Content to analyze:
{content[:8000]}  # Limit content to avoid token limits
"""

        try:
            response = self.claude_client.messages.create(
                model="claude-sonnet-4-20250514",
                max_tokens=2000,
                messages=[{
                    "role": "user",
                    "content": prompt
                }]
            )
            
            # Parse JSON response
            response_text = response.content[0].text.strip()
            
            # Clean up response (remove any markdown formatting)
            if response_text.startswith('```json'):
                response_text = response_text.replace('```json', '').replace('```', '').strip()
            
            metrics_data = json.loads(response_text)
            return metrics_data
            
        except json.JSONDecodeError as e:
            print(f"❌ Error parsing Claude response as JSON: {e}")
            print(f"Raw response: {response_text[:200]}...")
            return None
        except Exception as e:
            print(f"❌ Error calling Claude API: {e}")
            return None

    def save_financial_metrics(self, metrics_data, company_id, email_update_id, source_info, session):
        """Save extracted financial metrics to database"""
        try:
            # Parse reporting date
            reporting_date = None
            if metrics_data.get('reporting_date') and metrics_data['reporting_date'] != 'N/A':
                try:
                    reporting_date = datetime.strptime(metrics_data['reporting_date'], '%Y-%m-%d')
                except:
                    reporting_date = None
            
            # Create financial metrics record
            financial_metrics = FinancialMetrics(
                company_id=company_id,
                email_update_id=email_update_id,
                reporting_period=metrics_data.get('reporting_period'),
                reporting_date=reporting_date,
                
                # Revenue metrics
                mrr=metrics_data.get('mrr'),
                arr=metrics_data.get('arr'),
                qrr=metrics_data.get('qrr'),
                total_revenue=metrics_data.get('total_revenue'),
                gross_revenue=metrics_data.get('gross_revenue'),
                net_revenue=metrics_data.get('net_revenue'),
                
                # Growth metrics
                mrr_growth=metrics_data.get('mrr_growth'),
                arr_growth=metrics_data.get('arr_growth'),
                revenue_growth_yoy=metrics_data.get('revenue_growth_yoy'),
                revenue_growth_mom=metrics_data.get('revenue_growth_mom'),
                
                # Financial health
                cash_balance=metrics_data.get('cash_balance'),
                net_burn=metrics_data.get('net_burn'),
                gross_burn=metrics_data.get('gross_burn'),
                runway_months=metrics_data.get('runway_months'),
                
                # Profitability
                gross_margin=metrics_data.get('gross_margin'),
                ebitda=metrics_data.get('ebitda'),
                ebitda_margin=metrics_data.get('ebitda_margin'),
                net_income=metrics_data.get('net_income'),
                
                # Customer metrics
                customer_count=metrics_data.get('customer_count'),
                new_customers=metrics_data.get('new_customers'),
                churn_rate=metrics_data.get('churn_rate'),
                ltv=metrics_data.get('ltv'),
                cac=metrics_data.get('cac'),
                
                # Operational metrics
                team_size=metrics_data.get('team_size'),
                bookings=metrics_data.get('bookings'),
                pipeline=metrics_data.get('pipeline'),
                
                # Context
                key_highlights=metrics_data.get('key_highlights'),
                key_challenges=metrics_data.get('key_challenges'),
                funding_status=metrics_data.get('funding_status'),
                
                # Source info
                source_type=source_info.get('type', 'email'),
                source_file=source_info.get('filename'),
                extraction_confidence=metrics_data.get('extraction_confidence', 'medium'),
                extraction_notes=f"Extracted via Claude AI from {source_info.get('type', 'email')}"
            )
            
            session.add(financial_metrics)
            session.flush()  # Get the ID
            
            return financial_metrics
            
        except Exception as e:
            print(f"❌ Error saving financial metrics: {e}")
            return None

    def process_email_update(self, email_update_id):
        """Process a single email update and its attachments"""
        session = SessionLocal()
        
        try:
            # Get email update
            email_update = session.query(EmailUpdate).filter(
                EmailUpdate.id == email_update_id
            ).first()
            
            if not email_update:
                print(f"❌ Email update {email_update_id} not found")
                return False
            
            company_name = email_update.company.name
            print(f"\n💰 Processing financial metrics for: {company_name}")
            print(f"📧 Email: {email_update.subject[:60]}...")
            
            # Check if already processed
            existing_metrics = session.query(FinancialMetrics).filter(
                FinancialMetrics.email_update_id == email_update_id
            ).first()
            
            if existing_metrics:
                print(f"✅ Metrics already extracted for this email")
                return True
            
            processed_count = 0
            
            # Process email content first
            if email_update.body:
                print(f"📝 Analyzing email content...")
                
                source_info = {
                    'type': 'email',
                    'filename': None
                }
                
                metrics_data = self.analyze_content_with_claude(
                    email_update.body, 
                    company_name,
                    f"Email: {email_update.subject}"
                )
                
                if metrics_data:
                    financial_metrics = self.save_financial_metrics(
                        metrics_data, 
                        email_update.company_id,
                        email_update_id,
                        source_info,
                        session
                    )
                    
                    if financial_metrics:
                        processed_count += 1
                        print(f"✅ Extracted metrics from email content")
                        self.display_extracted_metrics(metrics_data)
            
            # Process attachments
            if email_update.attachments:
                print(f"📎 Processing {len(email_update.attachments)} attachments...")
                
                for attachment in email_update.attachments:
                    print(f"\n📄 Processing: {attachment.filename}")
                    
                    # Extract content from attachment
                    content = self.extract_attachment_content(attachment)
                    
                    if content:
                        source_info = {
                            'type': 'attachment',
                            'filename': attachment.filename
                        }
                        
                        metrics_data = self.analyze_content_with_claude(
                            content,
                            company_name,
                            f"Attachment: {attachment.filename}"
                        )
                        
                        if metrics_data:
                            financial_metrics = self.save_financial_metrics(
                                metrics_data,
                                email_update.company_id,
                                email_update_id,
                                source_info,
                                session
                            )
                            
                            if financial_metrics:
                                processed_count += 1
                                print(f"✅ Extracted metrics from {attachment.filename}")
                                self.display_extracted_metrics(metrics_data)
                        else:
                            print(f"❌ No metrics extracted from {attachment.filename}")
                    else:
                        print(f"❌ Could not extract content from {attachment.filename}")
            
            session.commit()
            
            print(f"\n📊 Summary: Extracted {processed_count} financial metric records")
            return processed_count > 0
            
        except Exception as e:
            print(f"❌ Error processing email update: {e}")
            session.rollback()
            return False
            
        finally:
            session.close()

    def display_extracted_metrics(self, metrics_data):
        """Display extracted metrics in a readable format"""
        print(f"   📅 Period: {metrics_data.get('reporting_period', 'N/A')}")
        
        # Show key financial metrics
        key_metrics = []
        if metrics_data.get('arr') and metrics_data['arr'] != 'N/A':
            key_metrics.append(f"ARR: {metrics_data['arr']}")
        if metrics_data.get('mrr') and metrics_data['mrr'] != 'N/A':
            key_metrics.append(f"MRR: {metrics_data['mrr']}")
        if metrics_data.get('cash_balance') and metrics_data['cash_balance'] != 'N/A':
            key_metrics.append(f"Cash: {metrics_data['cash_balance']}")
        if metrics_data.get('runway_months') and metrics_data['runway_months'] != 'N/A':
            key_metrics.append(f"Runway: {metrics_data['runway_months']}")
        
        if key_metrics:
            print(f"   💰 Key Metrics: {' | '.join(key_metrics)}")
        
        print(f"   🎯 Confidence: {metrics_data.get('extraction_confidence', 'medium')}")

    def process_recent_emails(self, days_back=30):
        """Process recent email updates for financial metrics extraction"""
        session = SessionLocal()
        
        try:
            # Get recent email updates that haven't been processed for metrics
            cutoff_date = datetime.now() - timedelta(days=days_back)
            
            recent_emails = session.query(EmailUpdate).filter(
                EmailUpdate.date >= cutoff_date
            ).outerjoin(FinancialMetrics).filter(
                FinancialMetrics.id.is_(None)  # Not yet processed
            ).order_by(EmailUpdate.date.desc()).all()
            
            print(f"💰 Financial Metrics Extraction")
            print(f"📅 Processing emails from last {days_back} days")
            print(f"📧 Found {len(recent_emails)} unprocessed emails")
            print("=" * 60)
            
            if not recent_emails:
                print("✅ All recent emails already processed for metrics")
                return
            
            processed_count = 0
            
            for email_update in recent_emails:
                if self.process_email_update(email_update.id):
                    processed_count += 1
            
            print(f"\n📊 Processing Complete!")
            print(f"✅ Successfully processed: {processed_count}/{len(recent_emails)} emails")
            
        except Exception as e:
            print(f"❌ Error processing recent emails: {e}")
            
        finally:
            session.close()


def main():
    """Main function for testing"""
    try:
        extractor = FinancialMetricsExtractor()
        extractor.process_recent_emails(days_back=30)
        
    except Exception as e:
        print(f"❌ Error: {e}")


if __name__ == "__main__":
    main() 